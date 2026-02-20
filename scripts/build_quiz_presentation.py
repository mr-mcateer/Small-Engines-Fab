#!/usr/bin/env python3
"""
Build "Anatomy of the Small Engine" Quiz Review Presentation
Apple minimalist style — clean, white, one idea per slide.
Includes "Think About It" coaching boxes and full speaker notes.

Also outputs Canvas answer-comment JSON for updating quiz feedback.

All wrong answers verified against Canvas quiz 41691 on 2026-02-18.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import json, os, hashlib, random

# ── Colors (Apple minimalist palette) ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
BLUE = RGBColor(0x00, 0x7A, 0xFF)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xFE)
GREEN = RGBColor(0x34, 0xC7, 0x59)
GRAYED = RGBColor(0xBB, 0xBB, 0xBB)
SECTION_BG = RGBColor(0xF5, 0xF5, 0xF7)

FONT = 'Helvetica Neue'
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Sections ──
SECTIONS = [
    {'title': 'External Components', 'sub': 'The parts you see on the outside of the engine',
     'qs': ['Q1','Q2','Q3','Q4','Q5','Q6'], 'img': 'recoil_starter.png'},
    {'title': 'Fuel & Air System', 'sub': 'Getting the right mix into the combustion chamber',
     'qs': ['Q7'], 'img': 'air_filter_carb.png'},
    {'title': 'Valvetrain', 'sub': 'Controlling what goes in and what comes out',
     'qs': ['Q8','Q9','Q10','Q11','Q12','Q19','Q20','Q28'], 'img': 'valvetrain.png'},
    {'title': 'Cylinder & Piston Assembly', 'sub': 'Where the magic happens — converting combustion to motion',
     'qs': ['Q13','Q14','Q15','Q16','Q17','Q18'], 'img': 'cylinder_piston.png'},
    {'title': 'Engine Systems', 'sub': 'Lubrication, speed control, energy storage, and power output',
     'qs': ['Q21','Q22','Q23','Q24','Q25'], 'img': 'flywheel.png'},
    {'title': 'Four-Stroke Cycle & Classification', 'sub': 'How it all comes together',
     'qs': ['Q26','Q27','Q29','Q30'], 'img': 'four_stroke_cycle.png'},
]

# ── All 30 Questions — wrong answers EXACTLY as they appear in Canvas ──
Q = {
    'Q1': {
        'text': 'What is the primary function of the recoil starter on a small engine?',
        'correct': 'Cranks the engine to start it',
        'wrong': ['Regulates engine speed', 'Cools the cylinder', 'Lubricates the crankshaft'],
        'think': 'What happens when you pull the starter cord? You\'re physically spinning the flywheel and crankshaft to get the first combustion cycle going.',
        'speaker': 'The recoil starter is the pull-cord mechanism mounted on top of the engine. When you pull it, the rope spins the flywheel, which turns the crankshaft, which moves the piston through its first cycles until the engine fires. Students should remember the physical sensation of pulling it during teardown — that resistance they felt was compression.',
        'wf': {
            'Regulates engine speed': 'That\'s the governor. The recoil starter is the pull cord — it physically spins the crankshaft to get the engine turning for the first time.',
            'Cools the cylinder': 'Cooling is handled by the shroud and fins. Think about what you grab and pull to start the engine.',
            'Lubricates the crankshaft': 'The oil slinger handles lubrication. The recoil starter is the pull-start mechanism — it cranks the engine.'
        }
    },
    'Q2': {
        'text': 'What component covers the top of the engine and directs airflow for cooling?',
        'correct': 'Shroud',
        'wrong': ['Crankcase', 'Valve cover', 'Muffler'],
        'think': 'Picture the engine from above. What\'s the big plastic/metal cover that sits over everything and has vents? It channels air from the flywheel fan across the cooling fins.',
        'speaker': 'The shroud is the large cover on top of the engine that works with the flywheel fan to direct air across the cylinder fins. Without it, the engine would overheat. During teardown, this was the first thing students removed. It\'s air-cooled — no radiator — so the shroud is critical.',
        'wf': {
            'Crankcase': 'The crankcase is the BOTTOM housing that holds the crankshaft and camshaft. Think about what covers the TOP and channels airflow.',
            'Valve cover': 'The valve cover is small and only covers the rocker arm area. The shroud is the large cover that directs cooling air over the entire engine.',
            'Muffler': 'The muffler handles exhaust gases and noise. The shroud is the cover on top that works with the flywheel fan to cool the engine.'
        }
    },
    'Q3': {
        'text': 'What does the fuel tank store?',
        'correct': 'Gasoline',
        'wrong': ['Engine oil', 'Coolant', 'Hydraulic fluid'],
        'think': 'The Predator 212 is a gasoline engine with a spark plug. Oil goes in the crankcase. This engine is air-cooled (no coolant). Hydraulic fluid is for hydraulic systems, not engines.',
        'speaker': 'Straightforward question that reinforces fundamentals. The Predator 212 is a spark-ignition gasoline engine. Engine oil goes in the crankcase through the oil fill cap — it\'s a completely separate system. This engine is air-cooled so there\'s no coolant. Hydraulic fluid is used in hydraulic systems like log splitters, not in the engine itself.',
        'wf': {
            'Engine oil': 'Engine oil goes in the CRANKCASE through the oil fill cap. The fuel tank is a separate container that holds gasoline for combustion.',
            'Coolant': 'This engine is AIR-COOLED — no coolant needed! The shroud and cooling fins handle heat dissipation. The fuel tank stores gasoline.',
            'Hydraulic fluid': 'Hydraulic fluid is for hydraulic systems (like log splitters or presses), not for the engine itself. The fuel tank stores gasoline.'
        }
    },
    'Q4': {
        'text': 'What is the purpose of the air filter on a small engine?',
        'correct': 'Prevents dirt and debris from entering the carburetor',
        'wrong': ['Filters the exhaust gases', 'Cleans the engine oil', 'Cools the incoming air'],
        'think': 'Follow the air path: outside air → air filter → carburetor → intake valve → combustion chamber. The filter is the first line of defense on the intake side.',
        'speaker': 'The air filter sits on the intake side before the carburetor. If dirt gets past it, it goes straight into the carburetor jets and then into the cylinder — causing wear and damage. During teardown, students saw how dirty these can get. A clogged air filter restricts airflow and makes the engine run rich.',
        'wf': {
            'Filters the exhaust gases': 'Exhaust goes OUT through the muffler. The air filter works on the INTAKE side — it cleans air BEFORE it enters the carburetor.',
            'Cleans the engine oil': 'Some engines have a separate oil filter, but the air filter only cleans incoming air. Follow the air path: outside → filter → carb → cylinder.',
            'Cools the incoming air': 'Air cooling comes from the shroud and fins. The air filter\'s only job is trapping dirt and debris before air reaches the carb.'
        }
    },
    'Q5': {
        'text': 'The spark plug is responsible for:',
        'correct': 'Igniting the air-fuel mixture in the combustion chamber',
        'wrong': ['Pumping fuel into the cylinder', 'Opening the exhaust valve', 'Compressing the piston rings'],
        'think': 'Follow the ignition circuit: magneto on flywheel → ignition coil → spark plug wire → spark plug → spark across the gap → ignites the mixture.',
        'speaker': 'The spark plug is the endpoint of the ignition system. The flywheel magnets spin past the ignition coil, generating a high-voltage pulse that travels down the spark plug wire and jumps the gap at the plug tip. This spark ignites the compressed air-fuel mixture during the power stroke. Students should remember the gap measurement from teardown.',
        'wf': {
            'Pumping fuel into the cylinder': 'Fuel delivery is the carburetor\'s job. The spark plug only creates the electrical spark needed for ignition.',
            'Opening the exhaust valve': 'Valves are opened mechanically by the camshaft and rocker arms. The spark plug sits in the head and just creates a spark.',
            'Compressing the piston rings': 'Piston rings are compressed by the cylinder wall — that\'s a mechanical fit. Follow the ignition circuit: coil → wire → plug → spark.'
        }
    },
    'Q6': {
        'text': 'What does the muffler do?',
        'correct': 'Reduces exhaust noise and directs exhaust gases away from the operator',
        'wrong': ['Increases engine horsepower', 'Filters incoming air', 'Stores excess fuel'],
        'think': 'Think about what comes OUT of the engine after combustion. Burnt gases exit through the exhaust port → muffler. It quiets the noise and points gases away safely.',
        'speaker': 'The muffler bolts to the exhaust port on the cylinder. After the power stroke, burnt gases are pushed out through the exhaust valve and into the muffler. The baffles inside reduce noise. It also directs hot gases away from the operator. Fun fact: removing or modifying the muffler is part of the Stage 1 upgrade on our drift trike build.',
        'wf': {
            'Increases engine horsepower': 'Stock mufflers actually create back-pressure that slightly reduces HP. The muffler\'s primary job is noise reduction and directing exhaust safely.',
            'Filters incoming air': 'That\'s the air filter on the intake side. The muffler is on the EXHAUST side — it handles gases leaving the engine.',
            'Stores excess fuel': 'The fuel tank stores fuel. The muffler only handles exhaust gases that have already been combusted.'
        }
    },
    'Q7': {
        'text': 'The carburetor\'s job is to:',
        'correct': 'Mix air and fuel in the correct ratio for combustion',
        'wrong': ['Ignite the fuel', 'Compress the air-fuel mixture', 'Lubricate the piston'],
        'think': 'The carburetor sits between the air filter and the intake valve. Its job is mixing — air and fuel in the right proportion (~14.7:1 by weight for gasoline).',
        'speaker': 'The carburetor uses the venturi effect — as air flows through a narrowing passage, it speeds up and creates low pressure, which draws fuel up through jets. The result is a fine mist of air and fuel in approximately a 14.7:1 ratio. Students adjusted the carburetor during Stage 1 mods. Too rich = black smoke, too lean = overheating.',
        'wf': {
            'Ignite the fuel': 'Ignition is the spark plug\'s job. The carburetor only MIXES air and fuel — it doesn\'t ignite anything.',
            'Compress the air-fuel mixture': 'Compression happens in the cylinder as the piston moves up. The carb just mixes air and fuel before they enter the cylinder.',
            'Lubricate the piston': 'Lubrication comes from the oil system (oil slinger). The carburetor only deals with air and fuel mixing.'
        }
    },
    'Q8': {
        'text': 'What is located under the valve cover?',
        'correct': 'Rocker arms and valve springs',
        'wrong': ['Piston and connecting rod', 'Crankshaft and camshaft', 'Oil slinger and governor'],
        'think': 'Remember pulling the valve cover off during teardown? Right underneath were the rocker arms sitting on top of the valve stems, with springs wrapped around them.',
        'speaker': 'When students removed the valve cover during teardown, they saw the rocker arms, valve springs, and the tops of the valve stems. The rocker arms pivot to push the valves open, and the springs pull them closed. This is a key visual memory from the hands-on work.',
        'wf': {
            'Piston and connecting rod': 'The piston and connecting rod are inside the cylinder and crankcase — much deeper in the engine. The valve cover is on top of the head.',
            'Crankshaft and camshaft': 'These are in the crankcase (bottom of the engine). The valve cover is on the cylinder HEAD and reveals the valvetrain components.',
            'Oil slinger and governor': 'These are in the crankcase (bottom of the engine). Under the valve cover on TOP you find the rocker arms and valve springs.'
        }
    },
    'Q9': {
        'text': 'What do the rocker arms do?',
        'correct': 'Transfer camshaft motion to open and close the valves',
        'wrong': ['Connect the piston to the crankshaft', 'Seal the combustion chamber', 'Regulate engine speed'],
        'think': 'Rocker arms are the mechanical link between the camshaft (below) and the valves (above). Cam lobe pushes up → pushrod → rocker arm pivots → valve opens.',
        'speaker': 'The rocker arms are lever mechanisms that transfer the camshaft\'s rotational motion into the linear push needed to open the valves. In the Predator 212 OHV design: cam lobe rotates → pushes tappet up → push rod goes up → rocker arm pivots → valve opens. When the cam lobe rotates away, the valve spring closes the valve.',
        'wf': {
            'Connect the piston to the crankshaft': 'That\'s the connecting rod\'s job. Rocker arms are up in the head, transferring cam motion to the valves.',
            'Seal the combustion chamber': 'The head gasket and piston rings handle sealing. Rocker arms are part of the valve train — they open and close the valves.',
            'Regulate engine speed': 'That\'s the governor\'s job. Rocker arms only control valve movement by transferring camshaft motion.'
        }
    },
    'Q10': {
        'text': 'The cylinder head is sealed to the cylinder by a:',
        'correct': 'Head gasket',
        'wrong': ['Piston ring', 'Valve spring', 'Crankcase cover'],
        'think': 'What sits BETWEEN the cylinder head and the cylinder block? A gasket — it seals the combustion chamber so no compression leaks out.',
        'speaker': 'The head gasket creates a seal between the cylinder head and the cylinder block. Without it, compression would leak out and the engine wouldn\'t run properly. During teardown, students saw this gasket when they removed the head. A blown head gasket is a common engine failure mode.',
        'wf': {
            'Piston ring': 'Piston rings seal the piston to the cylinder WALL (inside). The head gasket seals the head to the cylinder BLOCK (on top).',
            'Valve spring': 'Valve springs close the valves. The head gasket is the flat seal between the cylinder head and the block.',
            'Crankcase cover': 'The crankcase cover seals the bottom of the engine. The head gasket seals the TOP — between head and cylinder.'
        }
    },
    'Q11': {
        'text': 'How many valves does the Predator 212 have?',
        'correct': '2 (one intake, one exhaust)',
        'wrong': ['1 (intake only)', '4 (two intake, two exhaust)', '3 (two intake, one exhaust)'],
        'think': 'Single-cylinder OHV engine = one intake valve (lets fresh air-fuel in) and one exhaust valve (lets burnt gases out). Two valves total.',
        'speaker': 'The Predator 212 is a single-cylinder OHV engine with two valves: one intake and one exhaust. The intake valve opens during the intake stroke to let the air-fuel mixture in, and the exhaust valve opens during the exhaust stroke to let burnt gases out. Students saw both when they removed the head.',
        'wf': {
            '1 (intake only)': 'You need a way for exhaust to get OUT too. The Predator 212 has TWO valves — one intake and one exhaust.',
            '4 (two intake, two exhaust)': 'Four-valve designs exist in high-performance engines, but the Predator 212 is a simple single-cylinder with just 2 valves.',
            '3 (two intake, one exhaust)': 'Three-valve designs are rare. The Predator 212 keeps it simple: one intake valve + one exhaust valve = 2 total.'
        }
    },
    'Q12': {
        'text': 'Carbon deposits inside the combustion chamber are caused by:',
        'correct': 'Incomplete combustion of fuel and oil over time',
        'wrong': ['Too much airflow through the carburetor', 'Water in the fuel tank', 'A loose spark plug'],
        'think': 'Carbon is the residue from fuel and oil that didn\'t completely burn. Over time it builds up on the piston crown, valves, and chamber walls.',
        'speaker': 'Carbon deposits are the black crusty buildup students saw inside the combustion chamber during teardown. They come from fuel and oil that didn\'t fully combust. Running rich, oil burning past worn rings, and extended use all contribute. Cleaning carbon deposits is part of regular maintenance.',
        'wf': {
            'Too much airflow through the carburetor': 'Too much air makes the engine run LEAN, which actually produces less carbon. Carbon comes from incomplete combustion of fuel and oil.',
            'Water in the fuel tank': 'Water in fuel causes running problems and corrosion, not carbon deposits. Carbon is unburned fuel and oil residue.',
            'A loose spark plug': 'A loose spark plug causes compression loss and misfires. Carbon deposits come from the gradual buildup of incompletely burned fuel and oil.'
        }
    },
    'Q13': {
        'text': 'The cylinder (jug) is the part of the engine where:',
        'correct': 'The piston moves up and down during operation',
        'wrong': ['Fuel is stored before combustion', 'The camshaft is housed', 'The recoil starter attaches'],
        'think': 'The cylinder is the round bore — the tube where the piston slides up and down. It\'s the main working chamber of the engine.',
        'speaker': 'The cylinder (mechanics call it the "jug") is the precision-bored tube where the piston reciprocates. The cylinder wall must be smooth and round for the piston rings to seal properly. Students saw the bore condition when they pulled the piston out during teardown.',
        'wf': {
            'Fuel is stored before combustion': 'Fuel is stored in the fuel tank. The cylinder is where the piston moves up and down — it\'s the main working bore.',
            'The camshaft is housed': 'The camshaft lives in the crankcase (bottom of the engine). The cylinder is the bore above where the piston operates.',
            'The recoil starter attaches': 'The recoil starter mounts to the shroud/top of the engine. The cylinder is the internal bore where the piston moves.'
        }
    },
    'Q14': {
        'text': 'What is the purpose of piston rings?',
        'correct': 'Seal the gap between the piston and cylinder wall',
        'wrong': ['Connect the piston to the crankshaft', 'Hold the wrist pin in place', 'Open and close the valves'],
        'think': 'The piston is slightly smaller than the cylinder bore. The rings expand outward to seal that gap — keeping compression above and oil below.',
        'speaker': 'Piston rings ride in grooves on the piston and press outward against the cylinder wall. Compression rings seal combustion pressure above the piston, and the oil ring scrapes excess oil off the wall. Without good ring seal, you lose compression and burn oil. Students saw the ring grooves during teardown.',
        'wf': {
            'Connect the piston to the crankshaft': 'That\'s the wrist pin and connecting rod. Piston rings seal the gap between the piston and cylinder wall.',
            'Hold the wrist pin in place': 'Wrist pin retaining clips hold the wrist pin. Piston rings are completely separate — they seal against the cylinder wall.',
            'Open and close the valves': 'Valves are controlled by the camshaft and rocker arms. Piston rings seal the piston-to-cylinder gap.'
        }
    },
    'Q15': {
        'text': 'The wrist pin (piston pin) connects the:',
        'correct': 'Piston to the connecting rod',
        'wrong': ['Connecting rod to the crankshaft', 'Camshaft to the crankshaft', 'Valve to the rocker arm'],
        'think': 'Think about the mechanical chain: piston → [wrist pin] → connecting rod → crankshaft. The wrist pin is the pivot point between piston and rod.',
        'speaker': 'The wrist pin passes through the piston and the small end of the connecting rod, creating a pivot joint. This allows the connecting rod to swing as the piston moves straight up and down while the crankshaft rotates. Students pushed this pin out during teardown.',
        'wf': {
            'Connecting rod to the crankshaft': 'The rod journal on the crankshaft connects to the big end of the rod. The WRIST pin connects the piston to the small end of the rod.',
            'Camshaft to the crankshaft': 'Timing gears connect camshaft and crankshaft. The wrist pin links the piston to the connecting rod.',
            'Valve to the rocker arm': 'The valve stem sits directly in the rocker arm. The wrist pin is in the piston assembly — connecting piston to connecting rod.'
        }
    },
    'Q16': {
        'text': 'What does "bore condition" refer to when inspecting a cylinder?',
        'correct': 'The condition of the inner cylinder wall surface',
        'wrong': ['The size of the exhaust port', 'The depth of the combustion chamber', 'The number of piston rings installed'],
        'think': '"Bore" = the inside of the cylinder. "Condition" = is the surface smooth, scratched, scored, or worn? You\'re checking the cylinder wall quality.',
        'speaker': 'Bore condition describes the state of the cylinder wall surface. A healthy bore has fine crosshatch marks from honing that help retain oil. A damaged bore shows scoring, scratches, or glazing. During teardown, students looked inside the cylinder for these conditions. Bad bore = poor ring seal = lost compression.',
        'wf': {
            'The size of the exhaust port': '"Bore" refers to the cylinder interior, not the exhaust port. You\'re inspecting the wall surface for scoring, wear, or glazing.',
            'The depth of the combustion chamber': 'Combustion chamber volume is a different measurement. "Bore condition" means the quality of the inner cylinder wall surface.',
            'The number of piston rings installed': 'Ring count is a design spec. "Bore condition" = the state of the cylinder wall surface — smooth, scratched, scored, or worn.'
        }
    },
    'Q17': {
        'text': 'The crankshaft converts:',
        'correct': 'The up-and-down motion of the piston into rotational motion',
        'wrong': ['Rotational motion into up-and-down motion', 'Fuel into electricity', 'Air pressure into vacuum'],
        'think': 'Picture the piston going up and down. What shape converts linear motion to rotation? A crank — just like a bicycle pedal converts your up-down leg motion into wheel rotation.',
        'speaker': 'The crankshaft converts the piston\'s reciprocating motion into rotational motion that can drive equipment. Think of a bicycle pedal: your legs go up and down, but the crank converts that into wheel rotation. The connecting rod links the piston to the crankshaft\'s offset journal.',
        'wf': {
            'Rotational motion into up-and-down motion': 'You\'ve got it backwards! The piston goes up and down, and the crankshaft converts that INTO rotation. Think: bicycle pedal.',
            'Fuel into electricity': 'The engine converts fuel to mechanical energy, not electricity. The crankshaft specifically converts linear piston motion into rotation.',
            'Air pressure into vacuum': 'The crankshaft doesn\'t deal with air pressure. It converts the piston\'s up-and-down motion into useful rotational motion.'
        }
    },
    'Q18': {
        'text': 'The connecting rod links the:',
        'correct': 'Piston to the crankshaft',
        'wrong': ['Crankshaft to the camshaft', 'Valve to the rocker arm', 'Flywheel to the starter'],
        'think': 'Follow the power chain: combustion pushes piston down → connecting rod transfers that force → crankshaft rotates. The rod links piston to crank.',
        'speaker': 'The connecting rod is the mechanical link between the piston and the crankshaft. The small end connects to the piston via the wrist pin, and the big end wraps around the crankshaft journal. We\'re installing a billet rod (ARC 6254) on the drift trike for strength.',
        'wf': {
            'Crankshaft to the camshaft': 'Timing gears connect crankshaft to camshaft. The connecting rod links the PISTON to the CRANKSHAFT.',
            'Valve to the rocker arm': 'The valve stem sits in the rocker arm directly. The connecting rod is deep inside the engine linking piston to crankshaft.',
            'Flywheel to the starter': 'The starter engages the flywheel through a cup/pawl system. The connecting rod links piston to crankshaft.'
        }
    },
    'Q19': {
        'text': 'What does the camshaft do?',
        'correct': 'Controls the opening and closing timing of the intake and exhaust valves',
        'wrong': ['Spins the flywheel', 'Pumps oil through the engine', 'Ignites the air-fuel mixture'],
        'think': 'The camshaft has lobes (bumps) that push against the valve train at precise times. Each lobe controls when a valve opens and for how long.',
        'speaker': 'The camshaft has egg-shaped lobes that push against tappets at precisely timed intervals. As the camshaft rotates, each lobe opens a valve through the pushrod/rocker arm chain. The cam profile determines lift and duration. Cam upgrades are a common mod, though we\'re skipping that on our drift trike.',
        'wf': {
            'Spins the flywheel': 'The crankshaft spins the flywheel directly. The camshaft\'s job is controlling valve timing — when the intake and exhaust valves open and close.',
            'Pumps oil through the engine': 'The oil slinger handles lubrication in the Predator 212. The camshaft controls valve opening and closing timing.',
            'Ignites the air-fuel mixture': 'Ignition is the spark plug\'s job. The camshaft mechanically controls when the valves open and close.'
        }
    },
    'Q20': {
        'text': 'Timing marks on the crankshaft and camshaft gears ensure that:',
        'correct': 'The valves open and close at the correct time relative to piston position',
        'wrong': ['The engine runs at the correct RPM', 'The oil pump delivers the right pressure', 'The spark plug fires at maximum compression'],
        'think': 'The crankshaft and camshaft must be synchronized. If timing is off, a valve could be open when the piston comes up — catastrophic. The timing marks align them correctly.',
        'speaker': 'Timing marks are reference dots on the crankshaft gear and camshaft gear. When assembled, these marks must align so valve events happen at exactly the right piston positions. If timing is off, the piston could hit an open valve. Students should remember aligning these marks during reassembly.',
        'wf': {
            'The engine runs at the correct RPM': 'RPM is controlled by the governor and throttle. Timing marks synchronize the crankshaft and camshaft so valves open at the right time.',
            'The oil pump delivers the right pressure': 'The Predator 212 uses splash lubrication (oil slinger), not a pump. Timing marks ensure VALVE timing is synchronized with piston position.',
            'The spark plug fires at maximum compression': 'Spark timing is controlled by the flywheel magneto position. Timing marks on the GEARS ensure the VALVES are synchronized with piston position.'
        }
    },
    'Q21': {
        'text': 'The crankcase houses which of the following components?',
        'correct': 'Crankshaft, camshaft, and connecting rod',
        'wrong': ['Piston, cylinder, and spark plug', 'Carburetor, air filter, and fuel tank', 'Valve cover, rocker arms, and valves'],
        'think': 'The crankcase is the lower half of the engine block. What lives INSIDE the bottom? The rotating parts: crankshaft, camshaft, connecting rod.',
        'speaker': 'The crankcase is the lower engine housing containing the crankshaft, camshaft (connected by timing gears), and the lower portion of the connecting rod. It also holds the engine oil in its sump. During teardown, splitting the crankcase revealed all these components.',
        'wf': {
            'Piston, cylinder, and spark plug': 'The piston and cylinder are ABOVE the crankcase, and the spark plug is in the head. The crankcase houses the crankshaft, camshaft, and connecting rod.',
            'Carburetor, air filter, and fuel tank': 'All EXTERNAL components mounted on the outside. The crankcase is the internal lower housing with the crankshaft, camshaft, and connecting rod.',
            'Valve cover, rocker arms, and valves': 'These are in the UPPER part (the head area). The crankcase is the LOWER housing with the crankshaft, camshaft, and connecting rod.'
        }
    },
    'Q22': {
        'text': 'What is the function of the governor on a small engine?',
        'correct': 'Regulates engine speed under varying loads to prevent over-revving',
        'wrong': ['Controls the ignition timing', 'Pumps oil to the bearings', 'Engages the recoil starter'],
        'think': 'The governor is a speed-limiting system. When the engine speeds up, the governor reduces throttle. When it bogs down, the governor opens the throttle. It maintains steady RPM.',
        'speaker': 'The governor uses centrifugal weights to sense engine speed and automatically adjust the throttle. On our drift trike, we\'re removing the governor for direct throttle control — but that means the rider must manage RPM manually.',
        'wf': {
            'Controls the ignition timing': 'Ignition timing is set by the flywheel magneto position. The governor controls ENGINE SPEED by automatically adjusting the throttle.',
            'Pumps oil to the bearings': 'The Predator 212 uses an oil slinger (splash lubrication), not a pump. The governor regulates engine speed.',
            'Engages the recoil starter': 'The recoil starter is a manual pull mechanism. The governor is an automatic speed regulator that adjusts throttle based on load.'
        }
    },
    'Q23': {
        'text': 'The oil slinger in the Predator 212 lubricates the engine by:',
        'correct': 'Spinning and splashing oil from the sump onto internal components',
        'wrong': ['Pumping oil under high pressure through galleries', 'Dripping oil from the valve cover onto the piston', 'Mixing oil directly into the fuel'],
        'think': 'The Predator 212 uses splash lubrication, not a pressurized oil pump. A paddle on the camshaft gear dips into the oil sump and flings oil everywhere.',
        'speaker': 'The Predator 212 uses splash lubrication — a small paddle (oil slinger) attached to the camshaft gear dips into the oil sump as it rotates and flings oil onto internals. No oil pump or pressurized system. This is why keeping the correct oil level is critical.',
        'wf': {
            'Pumping oil under high pressure through galleries': 'That\'s a pressurized oil system (used in cars). The Predator 212 is simpler — a slinger SPLASHES oil from the sump.',
            'Dripping oil from the valve cover onto the piston': 'Oil doesn\'t drip down from above. The slinger at the BOTTOM spins and flings oil upward from the sump.',
            'Mixing oil directly into the fuel': 'That\'s a TWO-STROKE engine method. The Predator 212 is FOUR-STROKE — oil stays in the crankcase, splashed by the slinger.'
        }
    },
    'Q24': {
        'text': 'The flywheel\'s primary functions include:',
        'correct': 'Storing rotational energy to keep the crankshaft spinning between power strokes and generating a magnetic field for ignition',
        'wrong': ['Cooling the engine and filtering air', 'Compressing the fuel mixture', 'Connecting the carburetor to the intake valve'],
        'think': 'The flywheel is a heavy wheel on the crankshaft. Its mass stores kinetic energy AND it has magnets that spin past the ignition coil to generate spark.',
        'speaker': 'The flywheel serves dual purposes: mass stores rotational energy through the non-power strokes, and permanent magnets generate the ignition pulse. We\'re installing a billet flywheel (ARC 6695) — lighter, stronger, better balance at higher RPM.',
        'wf': {
            'Cooling the engine and filtering air': 'The flywheel fan helps with cooling, but its PRIMARY functions are energy storage and ignition generation. Air filtering is the air filter\'s job.',
            'Compressing the fuel mixture': 'The piston compresses the mixture in the cylinder. The flywheel stores rotational energy and generates the magnetic field for spark.',
            'Connecting the carburetor to the intake valve': 'The intake manifold connects carb to valve. The flywheel is on the crankshaft — it stores energy and generates ignition pulses.'
        }
    },
    'Q25': {
        'text': 'The PTO (Power Take-Off) side of the engine is where:',
        'correct': 'The output shaft exits the engine to drive equipment or a drivetrain',
        'wrong': ['The recoil starter is mounted', 'The air filter is located', 'The muffler attaches'],
        'think': 'PTO = Power Take-Off. It\'s where power LEAVES the engine. The output shaft sticks out from this side to connect to whatever the engine drives.',
        'speaker': 'PTO stands for Power Take-Off — the crankshaft extends out from this side to attach driven equipment. In our case, the centrifugal clutch and sprocket for the drift trike drivetrain. The PTO side is opposite the flywheel/starter side.',
        'wf': {
            'The recoil starter is mounted': 'The recoil starter is on the FLYWHEEL side (opposite). PTO = Power Take-Off = where the output shaft exits.',
            'The air filter is located': 'The air filter is on the carburetor side. PTO side is where the output shaft exits — where power leaves the engine.',
            'The muffler attaches': 'The muffler bolts to the exhaust port on the cylinder. PTO side is where the output shaft exits to drive equipment.'
        }
    },
    'Q26': {
        'text': 'What are the four strokes in a four-stroke engine cycle, in order?',
        'correct': 'Intake, Compression, Power, Exhaust',
        'wrong': ['Compression, Intake, Exhaust, Power', 'Power, Exhaust, Intake, Compression', 'Intake, Power, Compression, Exhaust'],
        'think': 'Suck, Squeeze, Bang, Blow — the mechanic\'s shorthand. Intake (suck), Compression (squeeze), Power (bang), Exhaust (blow).',
        'speaker': 'The four strokes: 1) INTAKE — piston down, intake valve open, mixture drawn in. 2) COMPRESSION — both valves closed, piston up, mixture squeezed. 3) POWER — spark fires, piston pushed down. 4) EXHAUST — piston up, exhaust valve open, burnt gases out. Mnemonic: Suck, Squeeze, Bang, Blow.',
        'wf': {
            'Compression, Intake, Exhaust, Power': 'Wrong order. You must take in the mixture FIRST (Intake), then Compress, then Power, then Exhaust. Suck, Squeeze, Bang, Blow.',
            'Power, Exhaust, Intake, Compression': 'Starts mid-cycle. Correct order: Intake → Compression → Power → Exhaust. Suck, Squeeze, Bang, Blow.',
            'Intake, Power, Compression, Exhaust': 'Can\'t have Power before Compression! The mixture must be compressed FIRST. Correct: Intake → Compression → Power → Exhaust.'
        }
    },
    'Q27': {
        'text': 'During the compression stroke, what is happening?',
        'correct': 'Both valves are closed and the piston moves upward, compressing the air-fuel mixture',
        'wrong': ['The intake valve is open and fuel is entering the cylinder', 'The exhaust valve is open and gases are leaving the cylinder', 'The spark plug fires and the piston moves downward'],
        'think': 'Compression = squeeze. Both valves MUST be closed (sealed chamber), and the piston moves UP to compress the mixture.',
        'speaker': 'During compression, both valves are closed creating a sealed chamber. The piston moves from BDC to TDC, compressing the mixture. This raises temperature and pressure for a powerful combustion event. The compression ratio on the Predator 212 is about 8.5:1.',
        'wf': {
            'The intake valve is open and fuel is entering the cylinder': 'That\'s the INTAKE stroke. During COMPRESSION, both valves are CLOSED so the mixture is sealed in and can be squeezed.',
            'The exhaust valve is open and gases are leaving the cylinder': 'That\'s the EXHAUST stroke. During compression, both valves are closed and the piston pushes upward.',
            'The spark plug fires and the piston moves downward': 'That\'s the POWER stroke. Compression comes BEFORE power — piston moves UP with both valves closed.'
        }
    },
    'Q28': {
        'text': 'Valve springs are responsible for:',
        'correct': 'Closing the valves after the camshaft lobe pushes them open',
        'wrong': ['Opening the valves during the intake stroke', 'Sealing the cylinder head to the block', 'Connecting the rocker arms to the push rods'],
        'think': 'The camshaft OPENS the valve (pushes it). But what CLOSES it? The spring! It pulls the valve back to its seat when the cam lobe rotates away.',
        'speaker': 'Valve springs provide the return force that closes valves. The camshaft opens them; the spring closes them. Spring tension is critical — too weak and the valve "floats" at high RPM. We\'re skipping stiffer valve springs on the drift trike since we won\'t hit high RPM.',
        'wf': {
            'Opening the valves during the intake stroke': 'The CAMSHAFT opens the valves via rocker arms. Springs do the OPPOSITE — they CLOSE the valves after the cam releases them.',
            'Sealing the cylinder head to the block': 'That\'s the head gasket. Valve springs wrap around valve stems and provide closing force.',
            'Connecting the rocker arms to the push rods': 'The pushrod sits in a socket in the rocker arm. Valve springs are on the valves themselves, providing closing force.'
        }
    },
    'Q29': {
        'text': 'The Predator 212 is classified as what type of engine?',
        'correct': 'Single-cylinder, four-stroke, air-cooled, overhead valve (OHV)',
        'wrong': ['Two-cylinder, two-stroke, liquid-cooled', 'Single-cylinder, two-stroke, air-cooled', 'V-twin, four-stroke, liquid-cooled'],
        'think': 'Count the cylinders (1), count the strokes (4), how is it cooled (air — fins + shroud), where are the valves (overhead = in the head).',
        'speaker': 'The Predator 212 classification: Single-cylinder, Four-stroke, Air-cooled, OHV. Knowing how to classify an engine by these characteristics is fundamental to small engine work.',
        'wf': {
            'Two-cylinder, two-stroke, liquid-cooled': 'Wrong on all counts. The Predator 212 has ONE cylinder, FOUR strokes, and is AIR-cooled.',
            'Single-cylinder, two-stroke, air-cooled': 'Close — but it\'s FOUR-stroke, not two-stroke. A two-stroke fires every revolution; this fires every OTHER revolution.',
            'V-twin, four-stroke, liquid-cooled': 'ONE cylinder (not V-twin) and AIR-cooled (not liquid). It is four-stroke and OHV though.'
        }
    },
    'Q30': {
        'text': 'What does OHV (Overhead Valve) mean?',
        'correct': 'The valves are located in the cylinder head above the piston, with the camshaft in the block driving them via pushrods or tappets',
        'wrong': ['The camshaft is located above the valves in the cylinder head', 'The valves are located on the side of the cylinder block', 'The engine has no valves and uses ports instead'],
        'think': 'OHV = Overhead Valve. Valves are OVERHEAD (in the head), but the camshaft stays DOWN in the block. Cam pushes up through tappets and pushrods to reach the valves above.',
        'speaker': 'OHV means valves in the cylinder head, camshaft in the block. Motion reaches valves through tappets, pushrods, and rocker arms. Different from OHC (cam also in head) or flathead (valves in block beside cylinder).',
        'wf': {
            'The camshaft is located above the valves in the cylinder head': 'That\'s OHC (Overhead CAM) or DOHC. In OHV, the VALVES are overhead but the CAMSHAFT stays in the block.',
            'The valves are located on the side of the cylinder block': 'That\'s a flathead or L-head design. OHV = Overhead Valve — valves are UP in the cylinder HEAD.',
            'The engine has no valves and uses ports instead': 'That\'s a two-stroke with port-style intake/exhaust. The Predator 212 is four-stroke OHV — valves in the head, cam in the block.'
        }
    },
}

# ── Helper functions ──

def set_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def text_box(slide, l, t, w, h, text, size=18, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return tb

def answer_line(slide, l, t, w, letter, text, size=17, color=CHARCOAL, bold=False, bullet_color=None):
    tb = slide.shapes.add_textbox(l, t, w, Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = f"{letter}.  "
    r1.font.size = Pt(size)
    r1.font.color.rgb = bullet_color or color
    r1.font.bold = True
    r1.font.name = FONT
    r2 = p.add_run()
    r2.text = text
    r2.font.size = Pt(size)
    r2.font.color.rgb = color
    r2.font.bold = bold
    r2.font.name = FONT
    return tb

def think_box(slide, l, t, w, h, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Think About It:  "
    r1.font.size = Pt(13)
    r1.font.color.rgb = BLUE
    r1.font.bold = True
    r1.font.name = FONT
    r2 = p.add_run()
    r2.text = text
    r2.font.size = Pt(13)
    r2.font.color.rgb = CHARCOAL
    r2.font.name = FONT
    return shape

def shuffle_answers(q_num, correct, wrong):
    """Deterministic shuffle so question/answer slides match."""
    all_a = [correct] + wrong
    seed = int(hashlib.md5(q_num.encode()).hexdigest()[:8], 16)
    random.Random(seed).shuffle(all_a)
    return all_a

# ── Slide builders ──

def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    text_box(s, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             'Anatomy of the Small Engine', size=44, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             'Quiz Review — Study Guide', size=28, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    text_box(s, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             'CVHS Engines & Fabrication  |  Predator 212', size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    # Try image
    img = os.path.join(IMG_DIR, 'engine_full.png')
    if os.path.exists(img):
        try: s.shapes.add_picture(img, Inches(4), Inches(4.5), Inches(5), Inches(2.7))
        except: pass
    s.notes_slide.notes_text_frame.text = (
        "Welcome to the quiz review and study guide. We'll go through every question, "
        "explain the correct answer, and give you a 'Think About It' framework so you understand "
        "the reasoning — not just the answer. Use this to study for retakes or deepen your understanding."
    )

def section_slide(prs, title, sub, img_name=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, SECTION_BG)
    text_box(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
             title, size=40, bold=True, align=PP_ALIGN.CENTER)
    # Accent line
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.5), Inches(2), Pt(3))
    sh.fill.solid(); sh.fill.fore_color.rgb = BLUE; sh.line.fill.background()
    text_box(s, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
             sub, size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if img_name:
        img = os.path.join(IMG_DIR, img_name)
        if os.path.exists(img):
            try: s.shapes.add_picture(img, Inches(3), Inches(4.5), Inches(7), Inches(2.7))
            except: pass

def question_slide(prs, qn, qd):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    text_box(s, Inches(0.8), Inches(0.5), Inches(1.5), Inches(0.5),
             f'Question {qn[1:]}', size=14, color=BLUE, bold=True)
    text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(1.2),
             qd['text'], size=26, bold=True)
    answers = shuffle_answers(qn, qd['correct'], qd['wrong'])
    letters = ['A','B','C','D']
    y = Inches(2.8)
    for i, ans in enumerate(answers):
        answer_line(s, Inches(1.2), y + Inches(i*0.8), Inches(10), letters[i], ans, size=18)
    s.notes_slide.notes_text_frame.text = f"Question {qn}: {qd['text']}\nLet students think before advancing."
    return answers, letters

def answer_slide(prs, qn, qd, answers, letters):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    text_box(s, Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
             f'Question {qn[1:]} — Answer', size=14, color=GREEN, bold=True)
    text_box(s, Inches(0.8), Inches(1.0), Inches(10), Inches(1.0),
             qd['text'], size=22, bold=True)
    y = Inches(2.2)
    for i, ans in enumerate(answers):
        is_c = ans == qd['correct']
        answer_line(s, Inches(1.2), y + Inches(i*0.65), Inches(10), letters[i], ans,
                    size=17, color=BLUE if is_c else GRAYED, bold=is_c,
                    bullet_color=GREEN if is_c else GRAYED)
    think_box(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.2), qd['think'])
    s.notes_slide.notes_text_frame.text = (
        f"ANSWER: {qd['correct']}\n\n{qd['speaker']}\n\nTHINK ABOUT IT: {qd['think']}"
    )

def wrapup_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    text_box(s, Inches(1), Inches(1.2), Inches(11), Inches(1),
             'Key Takeaways', size=40, bold=True, align=PP_ALIGN.CENTER)
    tips = [
        "Know the air-fuel path: filter → carb → intake valve → cylinder → exhaust valve → muffler",
        "Know the valve chain: cam lobe → tappet → pushrod → rocker arm → valve opens → spring closes",
        "Know the cycle: Intake → Compression → Power → Exhaust (Suck, Squeeze, Bang, Blow)",
        "Know your engine: Single-cylinder, four-stroke, air-cooled, OHV",
        "When in doubt — trace the path of air, fuel, force, or motion through the engine",
    ]
    for i, t in enumerate(tips):
        text_box(s, Inches(1.2), Inches(2.6+i*0.75), Inches(10.5), Inches(0.6),
                 f"→  {t}", size=16)
    s.notes_slide.notes_text_frame.text = (
        "Five mental models for any small engine question. Trace the path — air, fuel, force, motion. "
        "If you can follow the chain start to finish, you can figure out what each component does."
    )

def end_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, SECTION_BG)
    text_box(s, Inches(1), Inches(2.5), Inches(11), Inches(1),
             'Crescent Valley High School', size=36, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             'Engines & Fabrication Program', size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    text_box(s, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             'Mr. McAteer  |  2025–2026', size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Canvas feedback JSON ──

def save_feedback_json():
    feedback = {}
    for qn, qd in Q.items():
        feedback[qn] = {
            'question': qd['text'],
            'correct_comment': f"Correct! {qd['think']}",
            'wrong_comments': qd.get('wf', {})
        }
    path = os.path.join(SCRIPT_DIR, 'canvas_answer_feedback.json')
    with open(path, 'w') as f:
        json.dump(feedback, f, indent=2)
    print(f"  Canvas feedback JSON → {path}")

# ── Main ──

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(SCRIPT_DIR, 'images')
OUTPUT = os.path.join(os.path.dirname(SCRIPT_DIR), 'student_docs',
                      'Anatomy_Small_Engine_Quiz_Review.pptx')

def main():
    print("Building quiz review presentation...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs)
    print("  [+] Title")

    for sec in SECTIONS:
        section_slide(prs, sec['title'], sec['sub'], sec.get('img'))
        print(f"  [+] Section: {sec['title']}")
        for qn in sec['qs']:
            qd = Q[qn]
            answers, letters = question_slide(prs, qn, qd)
            answer_slide(prs, qn, qd, answers, letters)
            print(f"      {qn}: {qd['text'][:50]}...")

    wrapup_slide(prs)
    print("  [+] Wrap-up")
    end_slide(prs)
    print("  [+] End")

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"\n  Saved: {OUTPUT}")
    print(f"  Slides: {len(prs.slides)}")

    save_feedback_json()
    print("\nDone!")

if __name__ == '__main__':
    main()
